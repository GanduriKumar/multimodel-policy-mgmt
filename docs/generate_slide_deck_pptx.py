"""
Generate a PowerPoint slide deck for the Multimodel Policy Management solution.
Covers 12 delivery-leader questions from the slidegenerator prompt.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ─── Colour palette ───────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF9, 0xFA)
DARK_NAVY  = RGBColor(0x0D, 0x1B, 0x2A)   # near-black for headings
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)   # body text
ACCENT1    = RGBColor(0x16, 0x5B, 0xAA)   # primary blue
ACCENT2    = RGBColor(0x00, 0xB4, 0xD8)   # teal highlight
ACCENT3    = RGBColor(0xF7, 0x7F, 0x00)   # amber call-out
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)   # card background
BORDER_BLUE= RGBColor(0x16, 0x5B, 0xAA)
MUTED_GRAY = RGBColor(0x6C, 0x75, 0x7D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─── Low-level helpers ────────────────────────────────────────────────────────

def add_shape_fill(shape, color: RGBColor, transparency: int = 0):
    """Solid fill on any shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor = None,
             line_color: RGBColor = None, line_width_pt: float = 0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=16, bold=False, italic=False,
                 color: RGBColor = DARK_TEXT, align=PP_ALIGN.LEFT,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color: RGBColor = DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before_pt=4, bullet=False, indent_level=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before_pt)
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


# ─── Slide builders ───────────────────────────────────────────────────────────

def set_slide_background(slide, color: RGBColor = WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_banner(slide, title: str, subtitle: str = ""):
    """Dark navy top banner with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), fill=DARK_NAVY)
    # Accent stripe
    add_rect(slide, 0, Inches(1.35), SLIDE_W, Inches(0.07), fill=ACCENT2)
    # Title text
    add_text_box(slide, title,
                 Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.75),
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.5), Inches(0.88), Inches(11.5), Inches(0.4),
                     font_size=14, italic=True, color=ACCENT2, align=PP_ALIGN.LEFT)


def add_footer(slide, text: str = "Multimodel Policy Management  |  Confidential"):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=DARK_NAVY)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.22), Inches(12.7), Inches(0.28),
                 font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)


def add_section_number(slide, num: int):
    """Amber circle badge with slide number."""
    circle = slide.shapes.add_shape(9, Inches(12.5), Inches(0.1), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT3
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(num)
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def build_bullet_card(slide, items: list, left, top, width, height,
                      bg: RGBColor = LIGHT_BLUE, icon: str = "✦"):
    add_rect(slide, left, top, width, height, fill=bg,
             line_color=BORDER_BLUE, line_width_pt=0.75)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.12),
        width - Inches(0.3), height - Inches(0.24)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "Calibri"


def build_two_col(slide, left_items, right_items,
                  left_header="", right_header="",
                  top=Inches(1.55), row_h=Inches(2.3)):
    col_w = Inches(6.0)
    gap = Inches(0.33)
    # Left column
    if left_header:
        add_text_box(slide, left_header, Inches(0.4), top,
                     col_w, Inches(0.4),
                     font_size=14, bold=True, color=ACCENT1)
    build_bullet_card(slide, left_items,
                      Inches(0.4), top + Inches(0.42), col_w, row_h)
    # Right column
    rx = Inches(0.4) + col_w + gap
    if right_header:
        add_text_box(slide, right_header, rx, top,
                     col_w, Inches(0.4),
                     font_size=14, bold=True, color=ACCENT1)
    build_bullet_card(slide, right_items,
                      rx, top + Inches(0.42), col_w, row_h)


def build_three_col(slide, cols: list, headers: list,
                    top=Inches(1.55), row_h=Inches(2.5)):
    """cols = list of bullet lists; headers = column header strings"""
    n = len(cols)
    total_w = Inches(12.5)
    col_w = total_w / n - Inches(0.1)
    for i, (hdr, items) in enumerate(zip(headers, cols)):
        lx = Inches(0.4) + i * (col_w + Inches(0.1))
        if hdr:
            add_text_box(slide, hdr, lx, top, col_w, Inches(0.4),
                         font_size=13, bold=True, color=ACCENT1)
        build_bullet_card(slide, items, lx, top + Inches(0.42), col_w, row_h)


def add_highlight_box(slide, label: str, value: str,
                      left, top, width=Inches(3.8), height=Inches(1.1)):
    """Metric / callout card."""
    add_rect(slide, left, top, width, height, fill=ACCENT1)
    # value (large)
    add_text_box(slide, value,
                 left + Inches(0.1), top + Inches(0.05),
                 width - Inches(0.2), Inches(0.6),
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # label (small)
    add_text_box(slide, label,
                 left + Inches(0.1), top + Inches(0.65),
                 width - Inches(0.2), Inches(0.38),
                 font_size=12, color=ACCENT2, align=PP_ALIGN.CENTER)


# ─── Individual slides ────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, WHITE)

    # Full-width hero banner
    add_rect(slide, 0, 0, SLIDE_W, Inches(4.8), fill=DARK_NAVY)
    add_rect(slide, 0, Inches(4.8), SLIDE_W, Inches(0.1), fill=ACCENT2)

    # Decorative geometric accent
    add_rect(slide, Inches(10.5), 0, Inches(2.83), Inches(4.8),
             fill=RGBColor(0x11, 0x26, 0x3E))

    # Icon-style watermark
    add_text_box(slide, "🛡", Inches(10.6), Inches(0.6), Inches(2.4), Inches(2.4),
                 font_size=80, color=RGBColor(0x1E, 0x45, 0x78), align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, "Multimodel Policy Management",
                 Inches(0.7), Inches(0.9), Inches(9.6), Inches(1.4),
                 font_size=42, bold=True, color=WHITE)

    # Subtitle
    add_text_box(slide,
                 "A Policy-First AI Safety & Compliance Layer\nfor Enterprise Delivery Teams",
                 Inches(0.7), Inches(2.4), Inches(9.6), Inches(1.0),
                 font_size=20, italic=True, color=ACCENT2)

    # Tagline
    add_text_box(slide,
                 "Define once. Enforce everywhere. Audit always.",
                 Inches(0.7), Inches(3.5), Inches(9.0), Inches(0.6),
                 font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Bottom meta strip
    add_rect(slide, 0, Inches(5.6), SLIDE_W, Inches(1.9), fill=OFF_WHITE)
    meta = [
        ("📅", "March 2026"),
        ("👤", "KumarGN  |  Author"),
        ("📜", "MIT Licensed"),
        ("🔧", "FastAPI · React · SQLite"),
    ]
    for i, (icon, txt) in enumerate(meta):
        lx = Inches(0.5) + i * Inches(3.2)
        add_text_box(slide, f"{icon}  {txt}",
                     lx, Inches(5.75), Inches(3.0), Inches(0.5),
                     font_size=13, color=DARK_NAVY)

    add_footer(slide)


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "1. The Delivery Challenge We're Solving",
               "What operational pain do teams face when deploying AI today?")
    add_section_number(slide, 1)

    pain_points = [
        "No centralized way to enforce compliance rules across AI applications",
        "Policy logic is duplicated across every app — inconsistent & hard to audit",
        "PII, harmful content, and hallucinations slip through without detection",
        "Zero audit trail  →  impossible to prove compliance to regulators",
    ]
    frequency = [
        "Occurs in virtually EVERY enterprise AI engagement today",
        "Regulatory review failures delay go-live by weeks or months",
        "Incident response is slow: no traceability = long investigation times",
        "Multiplies with scale — each new LLM or app adds new exposure",
    ]

    build_two_col(slide, pain_points, frequency,
                  left_header="⚠️  Operational Pain Points",
                  right_header="📈  How Often & How Badly?",
                  top=Inches(1.58), row_h=Inches(2.6))

    # Risk callout
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.85),
             fill=RGBColor(0xFF, 0xF3, 0xCD),
             line_color=ACCENT3, line_width_pt=1.0)
    add_text_box(slide,
                 "⚡  Without a unified policy layer, each delivery team reinvents the wheel — "
                 "creating inconsistent guardrails, audit gaps, and hidden compliance debt.",
                 Inches(0.55), Inches(4.54), Inches(12.2), Inches(0.75),
                 font_size=13, italic=True, color=RGBColor(0x6E, 0x3D, 0x00))
    add_footer(slide)


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "2. Why Solving This Matters Right Now",
               "Urgency, risk, and the cost of inaction")
    add_section_number(slide, 2)

    metrics = [
        ("Regulatory Acts Active", "EU AI Act\n+ GDPR"),
        ("Avg. Cost of AI Compliance\nFailure", "$4.45M\nper incident"),
        ("LLM Adoption Growth\nYoY", "3×\nscale"),
    ]
    gap = Inches(0.25)
    card_w = (Inches(12.5) - 2 * gap) / 3
    for i, (lbl, val) in enumerate(metrics):
        lx = Inches(0.4) + i * (card_w + gap)
        add_highlight_box(slide, lbl, val, lx, Inches(1.65), card_w, Inches(1.35))

    risks = [
        "EU AI Act and sector regulations now mandate auditable AI decision trails",
        "Data breaches from unguarded LLM outputs cost millions in fines & remediation",
        "Reputational damage from leaked PII or harmful AI outputs is irreversible",
        "Without guardrails, scaling AI multiplies risk — not just capability",
        "Competitors already treating AI governance as a differentiator, not an afterthought",
    ]
    build_bullet_card(slide, risks, Inches(0.4), Inches(3.2), Inches(12.5), Inches(2.15),
                      bg=LIGHT_BLUE, icon="🔴")

    add_text_box(slide,
                 "The window to embed governance early is NOW — retrofitting is 5× more expensive.",
                 Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.45),
                 font_size=14, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "3. How the Solution Works",
               "Core idea, mechanics, and where it fits in the delivery lifecycle")
    add_section_number(slide, 3)

    # Flow diagram (text-based)
    flow_steps = [
        ("1️⃣", "User\nPrompt", ACCENT1),
        ("→", "", WHITE),
        ("2️⃣", "Pre-Check\n/api/protect", ACCENT1),
        ("→", "", WHITE),
        ("3️⃣", "LLM Call\n(any provider)", RGBColor(0x19, 0x76, 0xD2)),
        ("→", "", WHITE),
        ("4️⃣", "Post-Check\n/api/protect", ACCENT1),
        ("→", "", WHITE),
        ("5️⃣", "Safe\nResponse", RGBColor(0x2E, 0x7D, 0x32)),
    ]
    step_w = Inches(1.35)
    gap_w  = Inches(0.1)
    total  = len(flow_steps)
    start_x = Inches(0.4)
    for i, (icon, label, col) in enumerate(flow_steps):
        lx = start_x + i * (step_w + gap_w)
        if icon == "→":
            add_text_box(slide, "➤", lx + Inches(0.1), Inches(2.05),
                         Inches(0.5), Inches(0.5),
                         font_size=22, color=MUTED_GRAY, align=PP_ALIGN.CENTER)
        else:
            add_rect(slide, lx, Inches(1.72), step_w, Inches(0.95), fill=col)
            add_text_box(slide, f"{icon}\n{label}",
                         lx, Inches(1.72), step_w, Inches(0.95),
                         font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Ledger strip
    add_rect(slide, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.4),
             fill=RGBColor(0xE3, 0xF2, 0xFD))
    add_text_box(slide,
                 "📋  Every step logged to immutable Governance Ledger  —  full audit trail",
                 Inches(0.55), Inches(2.84), Inches(12.0), Inches(0.35),
                 font_size=12, italic=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    cols = [
        ["Versioned JSON policy documents",
         "Blocked terms, PII rules, source allow-lists",
         "Risk threshold per policy (0-100 score)",
         "Activate / rollback a version anytime"],
        ["Deterministic policy engine (testable)",
         "Heuristic risk scoring engine",
         "Groundedness check against evidence",
         "Response safety engine for LLM output"],
        ["Pre-check user input before LLM call",
         "Post-check LLM output before delivery",
         "One-call /api/protect-generate endpoint",
         "Works across Ollama, OpenAI, Vertex AI"],
    ]
    build_three_col(slide, cols,
                    ["📝  Policy Layer", "⚙️  Evaluation Engines", "🔄  Integration Modes"],
                    top=Inches(3.32), row_h=Inches(2.5))
    add_footer(slide)


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "4. Maturity & Reliability",
               "What stage is the solution at, and what evidence exists?")
    add_section_number(slide, 4)

    stages = ["Concept", "Prototype", "✅ Pilot / MVP", "Production"]
    stage_w = Inches(2.9)
    for i, s in enumerate(stages):
        is_current = "✅" in s
        col = ACCENT1 if is_current else RGBColor(0xDE, 0xE2, 0xE6)
        txt_col = WHITE if is_current else MUTED_GRAY
        lx = Inches(0.4) + i * (stage_w + Inches(0.2))
        add_rect(slide, lx, Inches(1.65), stage_w, Inches(0.7), fill=col)
        add_text_box(slide, s, lx, Inches(1.65), stage_w, Inches(0.7),
                     font_size=15, bold=is_current, color=txt_col, align=PP_ALIGN.CENTER)

    evidence = [
        "✔  Working REST API (FastAPI) with full OpenAPI documentation at /docs",
        "✔  Policy engine, risk engine, groundedness & safety engines — all independently tested",
        "✔  Test suite: unit, integration, API tests with pytest (Makefile targets: make test)",
        "✔  Tamper-evident governance ledger with SHA-256 hash chaining & verification tool",
        "✔  Compliance export service — machine-verifiable JSON + PDF-ready HTML bundles",
        "✔  React dashboard: create policies, run checks, browse audit, export reports",
        "✔  Multi-provider LLM gateway — Ollama & OpenAI confirmed; Vertex AI ready to wire",
    ]
    build_bullet_card(slide, evidence, Inches(0.4), Inches(2.55), Inches(12.5), Inches(3.2),
                      bg=LIGHT_BLUE, icon="")
    add_footer(slide)


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "5. Measurable Benefits",
               "Productivity gains, risk reduction, and business outcomes")
    add_section_number(slide, 5)

    gap = Inches(0.2)
    card_w = (Inches(12.5) - 2 * gap) / 3
    kpis = [
        ("Policy Reuse\nAcross Apps", "1 policy →\n∞ apps", ACCENT1),
        ("Audit Time\nReduction", "Manual hours\n→ 1-click export", RGBColor(0x00, 0x89, 0x7B)),
        ("Compliance Risk\nReduction", "Systematic\nguardrails", ACCENT3),
    ]
    for i, (lbl, val, col) in enumerate(kpis):
        lx = Inches(0.4) + i * (card_w + gap)
        add_rect(slide, lx, Inches(1.65), card_w, Inches(1.25), fill=col)
        add_text_box(slide, val, lx, Inches(1.68), card_w, Inches(0.75),
                     font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, lx, Inches(2.42), card_w, Inches(0.48),
                     font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    left_benefits = [
        "Same policy enforced across ALL LLM providers and applications",
        "Policy versioning enables safe evolution without breaking existing apps",
        "Deterministic engines → explainable decisions (not 'black box' safety)",
        "Risk score (0-100) gives teams a quantitative handle on each request",
    ]
    right_benefits = [
        "One-click compliance bundles reduce audit prep from days to minutes",
        "Hash-chained ledger satisfies 'immutable evidence' regulatory requirements",
        "Pre + post checks eliminate both unsafe inputs AND unsafe LLM outputs",
        "UI accessible to non-engineers — no code needed for policy management",
    ]
    build_two_col(slide, left_benefits, right_benefits,
                  left_header="🎯  Quality & Risk",
                  right_header="⏱  Efficiency & Compliance",
                  top=Inches(3.1), row_h=Inches(2.4))
    add_footer(slide)


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "6. Ease of Adoption & Deployment",
               "How quickly can delivery teams start realizing value?")
    add_section_number(slide, 6)

    timeline = [
        ("Day 1", "Clone repo, run backend (uvicorn) + frontend (npm dev). First policy live in < 1 hour."),
        ("Day 2-3", "Integrate the 'sandwich' pattern into existing app: 2 REST calls wrap any LLM call."),
        ("Week 1", "Configure policies for your domain; enable governance ledger & compliance exports."),
        ("Week 2+", "Roll out to additional apps & LLM providers. No policy rewrites needed."),
    ]
    t_top = Inches(1.72)
    row_h = Inches(1.05)
    for i, (phase, desc) in enumerate(timeline):
        ty = t_top + i * (row_h + Inches(0.08))
        # Phase badge
        add_rect(slide, Inches(0.4), ty, Inches(1.4), row_h,
                 fill=ACCENT1 if i == 0 else RGBColor(0xE3, 0xF2, 0xFD))
        add_text_box(slide, phase,
                     Inches(0.4), ty, Inches(1.4), row_h,
                     font_size=14, bold=True,
                     color=WHITE if i == 0 else ACCENT1,
                     align=PP_ALIGN.CENTER)
        # Description
        add_rect(slide, Inches(1.9), ty, Inches(11.0), row_h,
                 fill=OFF_WHITE, line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width_pt=0.5)
        add_text_box(slide, desc,
                     Inches(2.1), ty + Inches(0.12), Inches(10.7), row_h - Inches(0.2),
                     font_size=14, color=DARK_TEXT)

    add_text_box(slide,
                 "🚀  No agent installation. No proprietary SDK. Just REST APIs your team already knows.",
                 Inches(0.4), Inches(6.38), Inches(12.5), Inches(0.45),
                 font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "7. Dependencies & Prerequisites",
               "What does your team need to get started?")
    add_section_number(slide, 7)

    hard_deps = [
        "Python 3.10+  (backend runtime — widely available)",
        "Node.js 18+ & npm  (frontend build — standard in most delivery environments)",
        "Git  (version control — universal)",
        "SQLite  (default DB, zero config)  OR  any SQLAlchemy-compatible DB",
    ]
    soft_deps = [
        "OpenAI API key  (optional — only if using OpenAI as LLM provider)",
        "Ollama installed locally  (optional — free, local LLM alternative)",
        "Vertex AI GCP credentials  (optional — placeholder, ready to wire)",
        "HMAC secret env var  (optional — enables tamper-evident ledger)",
    ]
    build_two_col(slide, hard_deps, soft_deps,
                  left_header="🔒  Required (Hard) Dependencies",
                  right_header="🔓  Optional / Configurable",
                  top=Inches(1.6), row_h=Inches(2.5))

    skills = [
        "Basic REST API knowledge to call /api/protect from any language",
        "No special ML / AI skills required — policy authoring is JSON",
        "UI is designed for non-engineers: create & manage policies without code",
        "Existing DevOps pipelines can deploy with standard Python + Node.js tooling",
    ]
    build_bullet_card(slide, skills, Inches(0.4), Inches(4.35), Inches(12.5), Inches(1.5),
                      bg=RGBColor(0xE8, 0xF5, 0xE9), icon="👩‍💻")
    add_footer(slide)


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "8. Alternatives & Competitive Comparison",
               "How does this solution stack up against current approaches?")
    add_section_number(slide, 8)

    headers = ["🛠  Ad-hoc Custom Filters", "☁️  Vendor Cloud Guardrails", "✅  This Solution"]
    cols = [
        ["Duplicated code across apps",
         "No versioning or audit trail",
         "Hard to test & maintain",
         "Breaks with each new LLM",
         "No compliance reporting"],
        ["Vendor lock-in",
         "Limited policy customization",
         "Black-box decisions",
         "Expensive at scale",
         "Compliance evidence sparse"],
        ["One policy → all apps & LLMs",
         "Full version history & rollback",
         "Transparent, testable engines",
         "LLM-provider agnostic",
         "1-click compliance bundles"],
    ]

    col_w = Inches(4.0)
    gap   = Inches(0.12)
    colors = [RGBColor(0xFF, 0xEC, 0xB3),
              RGBColor(0xE3, 0xF2, 0xFD),
              RGBColor(0xE8, 0xF5, 0xE9)]
    h_colors = [ACCENT3, RGBColor(0x19, 0x76, 0xD2), RGBColor(0x2E, 0x7D, 0x32)]
    for i, (hdr, items, bg, hc) in enumerate(zip(headers, cols, colors, h_colors)):
        lx = Inches(0.4) + i * (col_w + gap)
        add_rect(slide, lx, Inches(1.65), col_w, Inches(0.48), fill=hc)
        add_text_box(slide, hdr, lx, Inches(1.65), col_w, Inches(0.48),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        build_bullet_card(slide, items, lx, Inches(2.13), col_w, Inches(3.5),
                          bg=bg, icon="•")

    add_text_box(slide,
                 "Strategic Advantage: Governance is baked in — not bolted on.",
                 Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.45),
                 font_size=14, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "9. Customer Differentiators & Adoption Enablers",
               "What will convince clients to choose this solution?")
    add_section_number(slide, 9)

    unique_caps = [
        "🏆  Policy-once, enforce-everywhere — no per-app policy rework",
        "🔍  Explainable decisions with human-readable reasons & risk scores",
        "📋  Tamper-evident ledger satisfies 'immutable audit' regulatory mandates",
        "🌐  Provider-agnostic: swap or run multiple LLMs without changing policy",
        "🖥  No-code policy management UI — empowers compliance & business teams",
    ]
    barriers = [
        ("'We already have filters'",
         "→  Show how those filters lack versioning, audit trails & cross-app consistency"),
        ("'We're locked into a vendor'",
         "→  Gateway abstraction means the vendor's LLM is still used — just governed"),
        ("'Governance slows delivery'",
         "→  Demonstrate Day-1 integration in < 1 hour; sandwich pattern is 2 REST calls"),
        ("'We can't trust open-source'",
         "→  MIT license, clean architecture, full test suite, verifiable compliance exports"),
    ]

    add_text_box(slide, "🌟  Unique Capabilities",
                 Inches(0.4), Inches(1.58), Inches(5.8), Inches(0.38),
                 font_size=14, bold=True, color=ACCENT1)
    build_bullet_card(slide, unique_caps, Inches(0.4), Inches(1.94), Inches(5.8), Inches(3.6),
                      bg=LIGHT_BLUE, icon="")

    add_text_box(slide, "🛡  Common Objections & How to Handle Them",
                 Inches(6.65), Inches(1.58), Inches(6.25), Inches(0.38),
                 font_size=14, bold=True, color=ACCENT1)
    ty = Inches(1.94)
    for barrier, response in barriers:
        add_rect(slide, Inches(6.65), ty, Inches(6.25), Inches(0.83),
                 fill=RGBColor(0xFF, 0xF8, 0xE1),
                 line_color=ACCENT3, line_width_pt=0.5)
        add_text_box(slide, barrier,
                     Inches(6.8), ty + Inches(0.04), Inches(6.0), Inches(0.36),
                     font_size=12, bold=True, color=RGBColor(0x6E, 0x3D, 0x00))
        add_text_box(slide, response,
                     Inches(6.8), ty + Inches(0.38), Inches(6.0), Inches(0.38),
                     font_size=11, italic=True, color=DARK_TEXT)
        ty += Inches(0.9)
    add_footer(slide)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "10. Steady-State Operational Cost",
               "What does it cost to run this solution in production?")
    add_section_number(slide, 10)

    cost_items = [
        ("Infrastructure", "SQLite (default) = $0.  Swap to PostgreSQL/MySQL on existing managed DB for minimal cost."),
        ("LLM Compute", "Ollama = free, runs on existing hardware.  OpenAI = pay-per-token per your current usage."),
        ("Backend Hosting", "A single Python FastAPI process.  Deployable on any cloud VM or container for ~$10-50/mo."),
        ("Frontend Hosting", "Static React build.  Serve from S3/Blob storage or existing CDN — near-zero cost."),
        ("Licensing", "MIT open-source — $0 licensing cost.  No per-seat or per-call fees."),
        ("Operational Effort", "Low.  Policy updates via UI (no code deployments).  Ledger is append-only (no maintenance)."),
    ]

    row_h = Inches(0.72)
    for i, (category, detail) in enumerate(cost_items):
        ty = Inches(1.65) + i * (row_h + Inches(0.06))
        bg = LIGHT_BLUE if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.4), ty, Inches(12.5), row_h,
                 fill=bg, line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width_pt=0.4)
        add_text_box(slide, category, Inches(0.55), ty + Inches(0.1),
                     Inches(2.0), row_h - Inches(0.1),
                     font_size=13, bold=True, color=ACCENT1)
        add_text_box(slide, detail, Inches(2.7), ty + Inches(0.1),
                     Inches(9.9), row_h - Inches(0.1),
                     font_size=13, color=DARK_TEXT)

    add_text_box(slide,
                 "💡  Total infra cost for a pilot: essentially $0 (local) to ~$50/mo (cloud VM + managed DB).",
                 Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.38),
                 font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_footer(slide)


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "11. Investment to Maintain & Evolve",
               "Long-term viability, ownership, and upgrade path")
    add_section_number(slide, 11)

    effort_areas = [
        ("🔧  Ongoing Maintenance",
         ["Policy updates: UI-driven, no engineering effort",
          "DB migrations: SQLAlchemy Alembic (scripted & versioned)",
          "Dependency upgrades: standard pip / npm upgrade cycle",
          "Ledger verification: automated via verify_ledger CLI tool"]),
        ("📈  Evolution & Extension",
         ["New LLM provider: implement ~50-line gateway adapter",
          "New policy rule type: extend PolicyDoc schema + engine",
          "New compliance format: extend ComplianceExportService",
          "Additional UI features: React + Vite component model"]),
        ("🏢  Organizational Ownership",
         ["One backend engineer part-time for ongoing maintenance",
          "Compliance / policy team owns policy authoring (UI, no code)",
          "Security team owns ledger HMAC secret rotation",
          "Clean architecture = easy onboarding for new engineers"]),
    ]

    col_w = Inches(3.9)
    gap = Inches(0.25)
    for i, (title, items) in enumerate(effort_areas):
        lx = Inches(0.4) + i * (col_w + gap)
        add_text_box(slide, title, lx, Inches(1.62), col_w, Inches(0.4),
                     font_size=13, bold=True, color=ACCENT1)
        build_bullet_card(slide, items, lx, Inches(2.05), col_w, Inches(3.55),
                          bg=LIGHT_BLUE, icon="→")

    add_rect(slide, Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.7),
             fill=RGBColor(0xE8, 0xF5, 0xE9),
             line_color=RGBColor(0x2E, 0x7D, 0x32), line_width_pt=0.75)
    add_text_box(slide,
                 "✅  The clean layered architecture (route → service → repo/engine) and full test suite "
                 "mean a new engineer can be productive within 1-2 days.",
                 Inches(0.6), Inches(5.88), Inches(12.1), Inches(0.58),
                 font_size=13, italic=True, color=RGBColor(0x1B, 0x5E, 0x20))
    add_footer(slide)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_banner(slide,
               "12. Transforming the Service Delivery Model",
               "The strategic takeaway for delivery leaders")
    add_section_number(slide, 12)

    before = [
        "Safety logic scattered across apps",
        "Compliance via manual checklists",
        "Audit = last-minute document scramble",
        "New LLM = rework safety from scratch",
        "Policy changes require code deployments",
        "Non-technical teams blocked from governance",
    ]
    after = [
        "One policy layer governs all AI touchpoints",
        "Compliance is continuous & automated",
        "Audit = 1-click export of verified evidence",
        "New LLM = plug into existing gateway",
        "Policy changes via UI in minutes, no code",
        "Business & compliance teams self-serve",
    ]

    add_text_box(slide, "❌  Before",
                 Inches(0.4), Inches(1.62), Inches(5.8), Inches(0.4),
                 font_size=14, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    build_bullet_card(slide, before, Inches(0.4), Inches(2.0), Inches(5.8), Inches(3.3),
                      bg=RGBColor(0xFF, 0xEB, 0xEB), icon="✗")

    add_text_box(slide, "✅  After",
                 Inches(6.65), Inches(1.62), Inches(5.8), Inches(0.4),
                 font_size=14, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
    build_bullet_card(slide, after, Inches(6.65), Inches(2.0), Inches(5.8), Inches(3.3),
                      bg=RGBColor(0xE8, 0xF5, 0xE9), icon="✓")

    # Arrow between columns
    add_text_box(slide, "➜",
                 Inches(6.1), Inches(3.2), Inches(0.55), Inches(0.7),
                 font_size=36, color=ACCENT1, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(1.72), fill=DARK_NAVY)
    add_text_box(slide,
                 "From reactive, app-by-app AI safety → to proactive, enterprise-wide AI governance.",
                 Inches(0.4), Inches(5.58), Inches(12.5), Inches(0.55),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "This is how trusted AI at scale is built — and it starts with a single policy definition.",
                 Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.45),
                 font_size=14, italic=True, color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(5.5), fill=DARK_NAVY)
    add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(0.1), fill=ACCENT2)
    add_rect(slide, Inches(10.3), 0, Inches(3.03), Inches(5.5),
             fill=RGBColor(0x11, 0x26, 0x3E))
    add_text_box(slide, "🛡", Inches(10.5), Inches(0.8), Inches(2.4), Inches(2.4),
                 font_size=80, color=RGBColor(0x1E, 0x45, 0x78), align=PP_ALIGN.CENTER)

    add_text_box(slide, "Ready to Get Started?",
                 Inches(0.7), Inches(0.9), Inches(9.3), Inches(1.0),
                 font_size=38, bold=True, color=WHITE)
    add_text_box(slide,
                 "Multimodel Policy Management is available today.\n"
                 "Define your first policy in under an hour.",
                 Inches(0.7), Inches(2.1), Inches(9.3), Inches(1.0),
                 font_size=18, italic=True, color=ACCENT2)

    links = [
        ("📖", "README & Quick Start", "github.com/your-org/multimodel-policy-mgmt"),
        ("📚", "User Guide",           "See UserGuide.md in the repo"),
        ("🔌", "API Docs",             "http://localhost:8000/docs  (after startup)"),
        ("📜", "License",              "MIT — free to use and modify"),
    ]
    for i, (icon, label, url) in enumerate(links):
        lx = Inches(0.7) + (i % 2) * Inches(4.6)
        ty = Inches(3.3) + (i // 2) * Inches(0.75)
        add_text_box(slide, f"{icon}  {label}",
                     lx, ty, Inches(2.0), Inches(0.5),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, url, lx + Inches(2.1), ty, Inches(2.4), Inches(0.5),
                     font_size=12, color=ACCENT2)

    add_rect(slide, 0, Inches(5.7), SLIDE_W, Inches(1.8), fill=OFF_WHITE)
    add_text_box(slide,
                 "\"Policy once. Enforce everywhere. Audit always.\"",
                 Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.6),
                 font_size=20, bold=True, italic=True,
                 color=ACCENT1, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "© 2026 KumarGN  |  MIT License  |  Built with FastAPI · React · SQLite",
                 Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.38),
                 font_size=11, color=MUTED_GRAY, align=PP_ALIGN.CENTER)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_closing(prs)

    out_path = r"c:\Users\kumar.gn\PycharmProjects\multimodel-policy-mgmt\docs\multimodel_policy_mgmt_slide_deck.pptx"
    prs.save(out_path)
    print(f"✅  Slide deck saved to:\n    {out_path}")
    print(f"    Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
