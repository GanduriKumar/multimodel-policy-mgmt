"""
Generate PowerPoint presentation from slide deck content.
Matches the HTML slide deck styling, colors, and typography.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color palette matching HTML (CSS variables)
COLORS = {
    'primary': RGBColor(13, 110, 253),      # #0d6efd
    'success': RGBColor(25, 135, 84),       # #198754
    'danger': RGBColor(220, 53, 69),        # #dc3545
    'warning': RGBColor(255, 193, 7),       # #ffc107
    'purple': RGBColor(111, 66, 193),       # #6f42c1
    'dark': RGBColor(33, 37, 41),           # #212529
    'light': RGBColor(248, 249, 250),       # #f8f9fa
    'gray': RGBColor(108, 117, 125),        # #6c757d
    'white': RGBColor(255, 255, 255),
    'gradient_blue': RGBColor(13, 71, 161), # #0d47a1
    'gradient_green': RGBColor(25, 135, 84),
    'orange': RGBColor(253, 126, 20),       # #fd7e14
}

def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 0: Title Slide
    create_title_slide(prs)
    
    # Slide 1: Problem Statement
    create_problem_slide(prs)
    
    # Slide 2: Key Components
    create_components_slide(prs)
    
    # Slide 2b: Key Components - Frontend & Reports
    create_frontend_slide(prs)
    
    # Slide 3: Competition & Differentiation
    create_competition_slide(prs)
    
    # Slide 4: Evolution Journey
    create_evolution_slide(prs)
    
    # Slide 5: Token Consumption
    create_token_slide(prs)
    
    # Slide 6: Summary
    create_summary_slide(prs)
    
    return prs


def create_title_slide(prs):
    """Slide 0: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Gradient background (approximation with solid color)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['gradient_blue']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "🛡️ Multimodel Policy Management"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.text = "Centralized AI governance: enforce regulations, manage risk, and maintain compliance across multiple LLM providers — from a single platform."
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    
    # Meta
    meta_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    tf = meta_box.text_frame
    tf.text = "KumarGN · 2025 · FastAPI + React + TypeScript\nMIT Licensed · Built for Enterprise AI Governance"
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
    
    # Slide number
    add_slide_number(slide, "0")


def create_problem_slide(prs):
    """Slide 1: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    add_section_title(slide, "1. Problem Statement")
    
    y_pos = 1.2
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = "The Enterprise AI Governance Gap"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.6
    
    # Problem description
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    tf = desc.text_frame
    tf.text = "As organizations adopt multiple LLM providers (OpenAI, Anthropic, open-source models), they face a critical challenge: how do you enforce consistent safety, compliance, and risk policies across all of them?"
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['dark']
    y_pos += 1.0
    
    # Four key problems in a 2x2 grid
    problems = [
        ("🔀", "Fragmented Policy Enforcement", "Each LLM provider has different content filters. Teams rewrite safety rules per app, per model — gaps and inconsistency."),
        ("📋", "Regulatory Compliance Burden", "EU AI Act, NIST AI RMF, NIST Privacy Framework demand auditable, traceable AI decisions. Manual compliance is unscalable."),
        ("🔍", "No Unified Audit Trail", "When an AI decision is challenged, organizations cannot easily prove what policy was active and why."),
        ("⚖️", "Risk Assessment Silos", "PII detection, intent classification, risk scoring scattered across disconnected tools."),
    ]
    
    card_width = 4.2
    card_height = 1.5
    x_positions = [0.5, 5.2]
    y_positions = [y_pos, y_pos + 1.8]
    
    for idx, (icon, title, desc_text) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        x = x_positions[col]
        y = y_positions[row]
        
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light']
        card.line.color.rgb = RGBColor(222, 226, 230)
        
        # Icon and title
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.1),
            Inches(card_width - 0.4), Inches(0.4)
        )
        tf = text_box.text_frame
        tf.text = f"{icon} {title}"
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.5),
            Inches(card_width - 0.4), Inches(0.9)
        )
        tf = desc_box.text_frame
        tf.text = desc_text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']
    
    # Footer and slide number
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "1")


def create_components_slide(prs):
    """Slide 2: Key Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_section_title(slide, "2. Key Components of the Solution")
    
    y_pos = 1.2
    
    # Architecture layers
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    tf.text = "Architecture — Clean Layered Design"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.5
    
    layers = [
        (COLORS['primary'], "Presentation Layer", "API Routes & React Dashboard"),
        (COLORS['purple'], "Application Layer", "Services & Orchestration"),
        (COLORS['success'], "Domain Layer", "Core Models & Engines"),
        (COLORS['orange'], "Infrastructure Layer", "Adapters & Persistence"),
    ]
    
    layer_height = 0.6
    for color, title, desc in layers:
        layer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(9), Inches(layer_height)
        )
        layer.fill.solid()
        layer.fill.fore_color.rgb = color
        layer.line.width = 0
        
        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.1),
            Inches(8.6), Inches(layer_height - 0.2)
        )
        tf = text_box.text_frame
        tf.text = f"{title}\n{desc}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(10)
            tf.paragraphs[1].font.color.rgb = COLORS['white']
        
        y_pos += layer_height + 0.1
    
    # Core components section
    y_pos += 0.3
    components_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3))
    tf = components_title.text_frame
    tf.text = "Core Components"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "2")


def create_frontend_slide(prs):
    """Slide 2b: Frontend & Reports"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_section_title(slide, "2. Key Components — Frontend & Reports")
    
    y_pos = 1.2
    
    # React Dashboard section
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    tf.text = "React Dashboard (TypeScript + Vite)"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.5
    
    # Dashboard features
    features = [
        ("📈", "Dashboard", "Live charts: decisions allow/deny over time, by policy, report generation"),
        ("⚡", "Protect Page", "Interactive policy testing with risk badges and evidence sources"),
        ("🔎", "Audit Page", "Browse requests/decisions with filters and compliance export"),
        ("📁", "Policies & Evidence", "CRUD for policies with JSON editor, version activation, side-by-side diff"),
    ]
    
    card_width = 4.2
    card_height = 1.3
    x_positions = [0.5, 5.2]
    y_positions = [y_pos, y_pos + 1.6]
    
    for idx, (icon, title, desc) in enumerate(features):
        row = idx // 2
        col = idx % 2
        x = x_positions[col]
        y = y_positions[row]
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light']
        card.line.color.rgb = RGBColor(222, 226, 230)
        
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.1),
            Inches(card_width - 0.4), Inches(1.1)
        )
        tf = text_box.text_frame
        tf.text = f"{icon} {title}\n{desc}"
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(10)
            tf.paragraphs[1].font.color.rgb = COLORS['dark']
    
    # Key design principle
    y_pos = 5.5
    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(y_pos),
        Inches(9), Inches(0.9)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(227, 242, 253)
    highlight.line.color.rgb = COLORS['primary']
    highlight.line.width = Pt(3)
    
    text_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(y_pos + 0.15),
        Inches(8.6), Inches(0.6)
    )
    tf = text_box.text_frame
    tf.text = "Key Design Principle: Swappable adapters behind Protocol ports — change LLM provider, database, or evidence store without touching business logic. Twelve-Factor compatible."
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark']
    
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "2b")


def create_competition_slide(prs):
    """Slide 3: Competition & Differentiation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_section_title(slide, "3. Competition & Differentiation")
    
    y_pos = 1.2
    
    # Key differentiators
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    tf.text = "Key Differentiators"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.5
    
    differentiators = [
        ("🏛️", "Regulation-Native", "Not 'add compliance later' — entire architecture designed around EU AI Act, NIST AI RMF, and privacy categories as first-class entities."),
        ("🔐", "Cryptographic Integrity", "Every compliance export has per-section SHA-256 hashes and root hash. Governance ledger uses hash-chaining with HMAC verification."),
        ("🧩", "Clean Architecture", "Strict layer separation: Domain has zero framework dependencies. All adapters swappable through Protocol ports and DI."),
    ]
    
    card_width = 9
    card_height = 1.3
    
    for icon, title, desc in differentiators:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light']
        card.line.color.rgb = RGBColor(222, 226, 230)
        
        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.15),
            Inches(card_width - 0.4), Inches(1.0)
        )
        tf = text_box.text_frame
        tf.text = f"{icon} {title}\n{desc}"
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(11)
            tf.paragraphs[1].font.color.rgb = COLORS['dark']
        
        y_pos += card_height + 0.2
    
    # Comparison highlights
    y_pos += 0.2
    comparison_title = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3))
    tf = comparison_title.text_frame
    tf.text = "Competitive Advantages vs. Alternatives"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.4
    
    advantages = [
        "✓ Multi-provider policy: Single policy works across any LLM",
        "✓ Tamper-evident audit: Hash-chained ledger with HMAC",
        "✓ Machine-verifiable exports: SHA-256 per-section + root hash",
        "✓ Human oversight workflows: Queues, SLA, approve/reject",
        "✓ Deterministic groundedness: No LLM needed for verification",
        "✓ Self-hosted & open source: MIT license, full control",
    ]
    
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.5))
    tf = text_box.text_frame
    for adv in advantages:
        p = tf.add_paragraph()
        p.text = adv
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['success']
        p.level = 0
    
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "3")


def create_evolution_slide(prs):
    """Slide 4: Evolution Journey"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_section_title(slide, "4. Feature Evolution Journey")
    
    y_pos = 1.2
    
    phases = [
        ("Phase 1", "Core Policy Engine", "Basic Protect API + Policy CRUD, risk scoring, request/decision logging", COLORS['primary']),
        ("Phase 2", "Risk Engine & Evidence", "Multi-signal risk, evidence management, groundedness engine, Protect & Generate API", COLORS['success']),
        ("Phase 3", "Compliance & Audit", "Compliance exports with SHA-256 hashes, governance ledger, framework-specific reports", COLORS['purple']),
        ("Phase 4", "Reporting & Visualization", "Decisions/Policy Changes reports, React Dashboard with Chart.js visualizations", COLORS['warning']),
        ("Phase 5", "Human Oversight & Maturity", "Review workflows with SLA, policy diff, privacy config, engineering constitution", COLORS['danger']),
    ]
    
    item_height = 1.1
    for phase, title, desc, color in phases:
        # Timeline indicator
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(y_pos + 0.15),
            Inches(0.25), Inches(0.25)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.color.rgb = COLORS['white']
        
        # Content box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), Inches(y_pos),
            Inches(8.5), Inches(item_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['light']
        box.line.color.rgb = RGBColor(222, 226, 230)
        
        # Text
        text_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(y_pos + 0.1),
            Inches(8.1), Inches(item_height - 0.2)
        )
        tf = text_box.text_frame
        tf.text = f"{phase} — {title}\n{desc}"
        tf.word_wrap = True
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(10)
            tf.paragraphs[1].font.color.rgb = COLORS['dark']
        
        y_pos += item_height + 0.15
    
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "4")


def create_token_slide(prs):
    """Slide 5: Token Consumption"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_section_title(slide, "5. Indicative Token Consumption")
    
    y_pos = 1.2
    
    # Introduction
    intro = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.6))
    tf = intro.text_frame
    tf.text = "Token usage varies by operation type. The system minimizes LLM token usage by performing most checks deterministically (policy matching, risk scoring, groundedness) and only calling LLMs for content generation."
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.8
    
    # Key insight callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(y_pos),
        Inches(9), Inches(0.7)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(255, 243, 205)
    callout.line.color.rgb = COLORS['warning']
    
    text_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(y_pos + 0.1),
        Inches(8.6), Inches(0.5)
    )
    tf = text_box.text_frame
    tf.text = "💡 Cost Optimization by Design: Policy evaluation, risk scoring, and groundedness checks are zero-token operations — they use deterministic algorithms, not LLM calls. Only content generation consumes tokens. >80% of requests consume zero LLM tokens."
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.9
    
    # Monthly estimates
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3))
    tf = subtitle.text_frame
    tf.text = "Monthly Estimates by Scale (GPT-4o-mini: $0.15/$0.60 per 1M input/output tokens)"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    y_pos += 0.4
    
    # Scale bars
    scales = [
        ("Small", "1K protect + 200 generate/month", "~$0.21", 15),
        ("Medium", "10K protect + 2K generate/month", "~$2.10", 40),
        ("Large", "100K protect + 20K generate/month", "~$21", 70),
        ("Enterprise", "1M protect + 200K generate/month", "~$210", 100),
    ]
    
    bar_height = 0.5
    colors_scale = [COLORS['success'], COLORS['primary'], COLORS['warning'], COLORS['danger']]
    
    for idx, (scale, desc, cost, width_pct) in enumerate(scales):
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos),
            Inches(6), Inches(0.25)
        )
        tf = label_box.text_frame
        tf.text = f"{scale} — {desc}"
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # Cost
        cost_box = slide.shapes.add_textbox(
            Inches(7), Inches(y_pos),
            Inches(2.5), Inches(0.25)
        )
        tf = cost_box.text_frame
        tf.text = cost
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.RIGHT
        
        y_pos += 0.3
        
        # Background bar
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(9), Inches(bar_height - 0.1)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(233, 236, 239)
        bg_bar.line.width = 0
        
        # Fill bar
        fill_width = 9 * (width_pct / 100)
        fill_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(fill_width), Inches(bar_height - 0.1)
        )
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = colors_scale[idx]
        fill_bar.line.width = 0
        
        y_pos += bar_height + 0.1
    
    add_footer(slide, "Multimodel Policy Management")
    add_slide_number(slide, "5")


def create_summary_slide(prs):
    """Slide 6: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['gradient_green']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Summary"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Key points
    y_pos = 2.5
    points = [
        ("🛡️", "One API, any LLM — consistent policy enforcement everywhere"),
        ("📜", "Versioned policies with full auditability and diff history"),
        ("🔐", "Cryptographic integrity — hash-chained ledger + per-section SHA-256"),
        ("📋", "Regulation-native — EU AI Act, NIST AI RMF, NIST Privacy built in"),
        ("💰", "Cost-efficient — 80%+ of requests use zero LLM tokens"),
        ("🏗️", "Clean architecture — swappable adapters, testable domain, DI throughout"),
    ]
    
    for icon, text in points:
        point_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(y_pos),
            Inches(7), Inches(0.4)
        )
        tf = point_box.text_frame
        tf.text = f"{icon}  {text}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']
        y_pos += 0.5
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    tf = thanks_box.text_frame
    tf.text = "Open Source · MIT Licensed · Built with FastAPI + React + TypeScript\n\nThank You"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(220, 220, 220)
    tf.paragraphs[-1].font.size = Pt(24)
    tf.paragraphs[-1].font.bold = True
    tf.paragraphs[-1].font.color.rgb = COLORS['white']
    
    add_slide_number(slide, "6")


def add_section_title(slide, text):
    """Add a section title with consistent styling."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Underline (approximated with a shape)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95),
        Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.width = 0


def add_footer(slide, text):
    """Add footer text."""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    tf = footer_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(173, 181, 189)


def add_slide_number(slide, number):
    """Add slide number."""
    num_box = slide.shapes.add_textbox(Inches(8.5), Inches(7.1), Inches(1), Inches(0.3))
    tf = num_box.text_frame
    tf.text = number
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['gray']


if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()
    output_path = "slide_deck.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
